"""
Document Processing Module
Provides parsing for multiple document formats with adaptive text chunking
Enhanced with semantic-aware chunking for improved retrieval accuracy
"""
import logging
from typing import List, Dict, Any, Optional
import os
from pathlib import Path
import json
import csv
import re

# Document parsers
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import markdown

logger = logging.getLogger(__name__)


class DocumentChunk:
    """Represents a chunk of document text with metadata"""
    
    def __init__(
        self,
        text: str,
        metadata: Dict[str, Any],
        chunk_index: int
    ):
        self.text = text
        self.metadata = metadata
        self.chunk_index = chunk_index
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return {
            'text': self.text,
            'metadata': self.metadata,
            'chunk_index': self.chunk_index
        }


class DocumentParser:
    """Parse various document formats and extract text"""
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """Parse PDF document"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    text += f"\n[Page {page_num + 1}]\n{page_text}"
            logger.info(f"Parsed PDF: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """Parse DOCX document"""
        try:
            doc = DocxDocument(file_path)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            logger.info(f"Parsed DOCX: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_txt(file_path: str) -> str:
        """Parse plain text document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            logger.info(f"Parsed TXT: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """Parse PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n[Slide {slide_num + 1}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            logger.info(f"Parsed PPTX: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_csv(file_path: str) -> str:
        """Parse CSV file"""
        try:
            text = ""
            with open(file_path, 'r', encoding='utf-8') as file:
                csv_reader = csv.reader(file)
                headers = next(csv_reader, None)
                if headers:
                    text += "Headers: " + ", ".join(headers) + "\n\n"
                
                for row_num, row in enumerate(csv_reader):
                    text += f"Row {row_num + 1}: " + " | ".join(row) + "\n"
            
            logger.info(f"Parsed CSV: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing CSV {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_json(file_path: str) -> str:
        """Parse JSON file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Convert JSON to readable text
            text = json.dumps(data, indent=2)
            logger.info(f"Parsed JSON: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing JSON {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_markdown(file_path: str) -> str:
        """Parse Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_text = file.read()
            
            # Keep original markdown but also convert to HTML for better structure
            html = markdown.markdown(md_text)
            # For simplicity, we'll use the original markdown text
            text = md_text
            
            logger.info(f"Parsed Markdown: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing Markdown {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_xlsx(file_path: str) -> str:
        """Parse Excel file"""
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n[Sheet: {sheet_name}]\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            
            logger.info(f"Parsed XLSX: {file_path}")
            return text
        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {e}")
            raise
    
    @staticmethod
    def parse_document(file_path: str) -> str:
        """
        Parse document based on file extension
        
        Args:
            file_path: Path to document file
            
        Returns:
            Extracted text content
        """
        file_extension = Path(file_path).suffix.lower()
        
        parsers = {
            '.pdf': DocumentParser.parse_pdf,
            '.docx': DocumentParser.parse_docx,
            '.doc': DocumentParser.parse_docx,
            '.txt': DocumentParser.parse_txt,
            '.pptx': DocumentParser.parse_pptx,
            '.ppt': DocumentParser.parse_pptx,
            '.csv': DocumentParser.parse_csv,
            '.json': DocumentParser.parse_json,
            '.md': DocumentParser.parse_markdown,
            '.markdown': DocumentParser.parse_markdown,
            '.xlsx': DocumentParser.parse_xlsx,
            '.xls': DocumentParser.parse_xlsx,
        }
        
        parser = parsers.get(file_extension)
        
        if parser is None:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        return parser(file_path)


class TextChunker:
    """Adaptive text chunking with semantic-aware boundaries"""
    
    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        min_chunk_size: int = 100,
        adaptive: bool = True
    ):
        """
        Initialize text chunker
        
        Args:
            chunk_size: Target size of each chunk in characters
            chunk_overlap: Number of overlapping characters between chunks
            min_chunk_size: Minimum size for a chunk to be valid
            adaptive: Use adaptive chunking strategy for better semantic coherence
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size
        self.adaptive = adaptive
        
        # Sentence boundary patterns
        self.sentence_endings = re.compile(r'[.!?]\s+')
        self.paragraph_split = re.compile(r'\n\s*\n')
        
        logger.info(f"Initialized text chunker (adaptive: {adaptive})")
    
    def chunk_text(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[DocumentChunk]:
        """
        Split text into overlapping chunks with adaptive strategy
        
        Args:
            text: Input text
            metadata: Optional metadata to attach to chunks
            
        Returns:
            List of DocumentChunk objects
        """
        if not text or len(text) < self.min_chunk_size:
            return []
        
        if self.adaptive:
            return self._adaptive_chunk(text, metadata)
        else:
            return self._simple_chunk(text, metadata)
    
    def _adaptive_chunk(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[DocumentChunk]:
        """
        Adaptive chunking that considers semantic boundaries
        Improves retrieval accuracy by maintaining context coherence
        
        Args:
            text: Input text
            metadata: Optional metadata
            
        Returns:
            List of DocumentChunk objects
        """
        chunks = []
        chunk_index = 0
        
        # First, split by paragraphs
        paragraphs = self.paragraph_split.split(text)
        
        current_chunk = ""
        current_start = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # If adding this paragraph would exceed chunk size
            if len(current_chunk) + len(para) > self.chunk_size:
                # If current chunk is not empty, save it
                if len(current_chunk) >= self.min_chunk_size:
                    chunk_metadata = metadata.copy() if metadata else {}
                    chunk_metadata['start_char'] = current_start
                    chunk_metadata['end_char'] = current_start + len(current_chunk)
                    chunk_metadata['chunking_strategy'] = 'adaptive'
                    
                    chunk = DocumentChunk(
                        text=current_chunk.strip(),
                        metadata=chunk_metadata,
                        chunk_index=chunk_index
                    )
                    chunks.append(chunk)
                    chunk_index += 1
                    
                    # Start new chunk with overlap
                    overlap_text = self._get_overlap_text(current_chunk, self.chunk_overlap)
                    current_start = current_start + len(current_chunk) - len(overlap_text)
                    current_chunk = overlap_text
                
                # If paragraph itself is too large, split by sentences
                if len(para) > self.chunk_size:
                    sentence_chunks = self._chunk_by_sentences(
                        para, 
                        current_start,
                        chunk_index,
                        metadata
                    )
                    chunks.extend(sentence_chunks)
                    chunk_index += len(sentence_chunks)
                    current_chunk = ""
                    current_start = current_start + len(para)
                else:
                    current_chunk = para
            else:
                # Add paragraph to current chunk
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
        
        # Don't forget the last chunk
        if len(current_chunk) >= self.min_chunk_size:
            chunk_metadata = metadata.copy() if metadata else {}
            chunk_metadata['start_char'] = current_start
            chunk_metadata['end_char'] = current_start + len(current_chunk)
            chunk_metadata['chunking_strategy'] = 'adaptive'
            
            chunk = DocumentChunk(
                text=current_chunk.strip(),
                metadata=chunk_metadata,
                chunk_index=chunk_index
            )
            chunks.append(chunk)
        
        logger.info(f"Created {len(chunks)} adaptive chunks from text of length {len(text)}")
        return chunks
    
    def _chunk_by_sentences(
        self,
        text: str,
        start_offset: int,
        start_index: int,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[DocumentChunk]:
        """
        Chunk large paragraph by sentences
        
        Args:
            text: Text to chunk
            start_offset: Character offset in original document
            start_index: Starting chunk index
            metadata: Optional metadata
            
        Returns:
            List of chunks
        """
        chunks = []
        sentences = self.sentence_endings.split(text)
        
        current_chunk = ""
        current_start = start_offset
        chunk_index = start_index
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            if len(current_chunk) + len(sentence) > self.chunk_size:
                if len(current_chunk) >= self.min_chunk_size:
                    chunk_metadata = metadata.copy() if metadata else {}
                    chunk_metadata['start_char'] = current_start
                    chunk_metadata['end_char'] = current_start + len(current_chunk)
                    chunk_metadata['chunking_strategy'] = 'sentence-based'
                    
                    chunk = DocumentChunk(
                        text=current_chunk.strip(),
                        metadata=chunk_metadata,
                        chunk_index=chunk_index
                    )
                    chunks.append(chunk)
                    chunk_index += 1
                    
                    # Overlap
                    overlap_text = self._get_overlap_text(current_chunk, self.chunk_overlap)
                    current_start = current_start + len(current_chunk) - len(overlap_text)
                    current_chunk = overlap_text + " " + sentence
                else:
                    current_chunk = sentence
            else:
                if current_chunk:
                    current_chunk += " " + sentence
                else:
                    current_chunk = sentence
        
        # Last chunk
        if len(current_chunk) >= self.min_chunk_size:
            chunk_metadata = metadata.copy() if metadata else {}
            chunk_metadata['start_char'] = current_start
            chunk_metadata['end_char'] = current_start + len(current_chunk)
            chunk_metadata['chunking_strategy'] = 'sentence-based'
            
            chunk = DocumentChunk(
                text=current_chunk.strip(),
                metadata=chunk_metadata,
                chunk_index=chunk_index
            )
            chunks.append(chunk)
        
        return chunks
    
    def _get_overlap_text(self, text: str, overlap_size: int) -> str:
        """
        Get overlap text by trying to maintain sentence boundary
        
        Args:
            text: Source text
            overlap_size: Desired overlap size
            
        Returns:
            Overlap text
        """
        if len(text) <= overlap_size:
            return text
        
        # Try to find sentence boundary in overlap region
        overlap_start = len(text) - overlap_size
        match = self.sentence_endings.search(text, overlap_start)
        
        if match:
            return text[match.end():]
        else:
            # Fallback to character-based overlap
            return text[-overlap_size:]
    
    def _simple_chunk(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[DocumentChunk]:
        """
        Simple character-based chunking (original implementation)
        
        Args:
            text: Input text
            metadata: Optional metadata
            
        Returns:
            List of DocumentChunk objects
        """
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(text):
            # Calculate end position
            end = start + self.chunk_size
            
            # If not at the end, try to break at sentence or word boundary
            if end < len(text):
                # Look for sentence boundary (., !, ?)
                for boundary in ['. ', '! ', '? ', '\n\n', '\n']:
                    boundary_pos = text.rfind(boundary, start, end)
                    if boundary_pos != -1:
                        end = boundary_pos + len(boundary)
                        break
                else:
                    # If no sentence boundary, look for word boundary
                    space_pos = text.rfind(' ', start, end)
                    if space_pos != -1:
                        end = space_pos + 1
            
            # Extract chunk
            chunk_text = text[start:end].strip()
            
            # Only add chunk if it meets minimum size
            if len(chunk_text) >= self.min_chunk_size:
                chunk_metadata = metadata.copy() if metadata else {}
                chunk_metadata['start_char'] = start
                chunk_metadata['end_char'] = end
                chunk_metadata['chunking_strategy'] = 'simple'
                
                chunk = DocumentChunk(
                    text=chunk_text,
                    metadata=chunk_metadata,
                    chunk_index=chunk_index
                )
                chunks.append(chunk)
                chunk_index += 1
            
            # Move to next chunk with overlap
            start = end - self.chunk_overlap
            
            # Avoid infinite loop
            if start <= 0 or end >= len(text):
                break
        
        logger.info(f"Created {len(chunks)} simple chunks from text of length {len(text)}")
        return chunks


class DocumentProcessor:
    """Process documents: parse and chunk with adaptive strategy"""
    
    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        min_chunk_size: int = 100,
        adaptive_chunking: bool = True,
        use_advanced_extractors: bool = False,
        tika_url: Optional[str] = None,
        docling_url: Optional[str] = None
    ):
        """
        Initialize document processor
        
        Args:
            chunk_size: Target chunk size
            chunk_overlap: Overlap between chunks
            min_chunk_size: Minimum chunk size
            adaptive_chunking: Use adaptive chunking strategy
            use_advanced_extractors: Try to use Tika/Docling if available
            tika_url: Apache Tika server URL
            docling_url: Docling server URL
        """
        self.parser = DocumentParser()
        self.chunker = TextChunker(
            chunk_size,
            chunk_overlap,
            min_chunk_size,
            adaptive=adaptive_chunking
        )
        
        # Initialize advanced extractors if requested
        self.use_advanced_extractors = use_advanced_extractors
        self.advanced_extractor = None
        
        if use_advanced_extractors:
            try:
                from .advanced_extractors import get_advanced_extractor
                self.advanced_extractor = get_advanced_extractor(
                    tika_url=tika_url,
                    docling_url=docling_url
                )
                if self.advanced_extractor.is_available():
                    logger.info("Advanced extractors enabled and available")
                else:
                    logger.info("Advanced extractors requested but not available, using default parsers")
            except Exception as e:
                logger.warning(f"Could not initialize advanced extractors: {e}")
    
    def process_document(
        self,
        file_path: str,
        document_id: Optional[str] = None
    ) -> tuple[str, List[DocumentChunk]]:
        """
        Process a document: parse and chunk
        
        Args:
            file_path: Path to document
            document_id: Optional document identifier
            
        Returns:
            Tuple of (full_text, chunks)
        """
        # Generate document ID if not provided
        if document_id is None:
            document_id = Path(file_path).stem
        
        # Try advanced extractors first if enabled
        text = None
        extraction_method = "default"
        
        if self.use_advanced_extractors and self.advanced_extractor:
            if self.advanced_extractor.is_available():
                text = self.advanced_extractor.extract_text(file_path)
                if text:
                    extraction_method = "advanced"
                    logger.info(f"Used advanced extractor for {file_path}")
        
        # Fallback to default parser
        if text is None:
            text = self.parser.parse_document(file_path)
            extraction_method = "default"
        
        # Create metadata
        metadata = {
            'document_id': document_id,
            'filename': Path(file_path).name,
            'file_extension': Path(file_path).suffix,
            'file_size': os.path.getsize(file_path),
            'extraction_method': extraction_method
        }
        
        # Chunk text
        chunks = self.chunker.chunk_text(text, metadata)
        
        logger.info(f"Processed document {file_path}: {len(chunks)} chunks (method: {extraction_method})")
        return text, chunks

